"""Deck generator template. Copy next to your HTML deck, rename to make_pptx.py,
edit the SLIDES section so every HTML slide has a matching PPTX slide, then run:
    python3 make_pptx.py
Produces Name.pptx (16:9, dark theme) with the same look as the HTML skeleton.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

W, H = 13.333, 7.5

BG     = RGBColor(0x0D, 0x10, 0x20)
PANEL  = RGBColor(0x17, 0x1B, 0x30)
PANEL2 = RGBColor(0x1E, 0x23, 0x40)
LINE   = RGBColor(0x2A, 0x30, 0x52)
INK    = RGBColor(0xF2, 0xF4, 0xFC)
MUTED  = RGBColor(0xA7, 0xAE, 0xCF)
VIOLET = RGBColor(0x6E, 0x7B, 0xFF)
GOLD   = RGBColor(0xFF, 0xC8, 0x57)
GREEN  = RGBColor(0x3D, 0xDC, 0x97)
CODEBG = RGBColor(0x0A, 0x0D, 0x1C)
CODE   = RGBColor(0xCD, 0xD6, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

MX = 0.7
CW = W - 2 * MX
CT = 1.62
FOOT_Y = 7.08
DECK_TITLE = "Deck Title"


def new_slide():
    s = prs.slides.add_slide(BLANK)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    b.fill.solid(); b.fill.fore_color.rgb = BG
    b.line.fill.background()
    b.shadow.inherit = False
    return s


def shape(s, kind, x, y, w, h):
    sh = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def rect(s, x, y, w, h, fill, line=None, ln_pt=0.75):
    r = shape(s, MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line:
        r.line.color.rgb = line; r.line.width = Pt(ln_pt)
    return r


def put_text(tf, text, size, bold=False, color=INK, font=FONT, align=None):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return p


def card(s, x, y, w, h, hd=None, body=None, lab=None, lab_fill=VIOLET):
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    c.fill.solid(); c.fill.fore_color.rgb = PANEL
    c.line.color.rgb = LINE; c.line.width = Pt(0.75)
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.11); tf.margin_bottom = Inches(0.08)
    off = 0.11
    if lab:
        lw = 0.16 + len(lab) * 0.085
        chip = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.16, y + 0.13, lw, 0.26)
        try:
            chip.adjustments[0] = 0.5
        except Exception:
            pass
        chip.fill.solid(); chip.fill.fore_color.rgb = lab_fill
        ct = chip.text_frame; ct.word_wrap = False
        ct.margin_left = Inches(0.04); ct.margin_right = Inches(0.04)
        ct.margin_top = 0; ct.margin_bottom = 0
        p = ct.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lab.upper()
        r.font.size = Pt(8.5); r.font.bold = True; r.font.name = FONT
        r.font.color.rgb = WHITE if lab_fill == VIOLET else (RGBColor(0x24, 0x1D, 0x00) if lab_fill == GOLD else RGBColor(0x03, 0x27, 0x16))
        off = 0.5
    first = True
    if hd:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(2)
        r = p.add_run(); r.text = hd
        r.font.size = Pt(15); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = INK
    if body:
        for b in body:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = b
            r.font.size = Pt(12); r.font.name = FONT
            r.font.color.rgb = RGBColor(0xBB, 0xC2, 0xE0)
    return y + h


def tag_and_title(s, tag, title, tsize=31):
    t = s.shapes.add_textbox(Inches(MX), Inches(0.28), Inches(CW), Inches(0.4))
    put_text(t.text_frame, tag.upper(), 11.5, True, GREEN)
    tt = s.shapes.add_textbox(Inches(MX), Inches(0.56), Inches(CW), Inches(0.95))
    put_text(tt.text_frame, title, tsize, True, INK)
    rect(s, MX, 1.52, 0.85, 0.045, GOLD)
    return CT


def bullets(s, x, y, w, items, size=14.5, gap=0.05):
    yy = y
    lh = size * 0.022 + 0.16
    for it in items:
        m = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, yy + 0.055, 0.11, 0.11)
        try:
            m.adjustments[0] = 0.3
        except Exception:
            pass
        m.fill.solid(); m.fill.fore_color.rgb = GREEN
        tb = s.shapes.add_textbox(Inches(x + 0.26), Inches(yy), Inches(w - 0.26), Inches(lh))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = it
        r.font.size = Pt(size); r.font.name = FONT; r.font.color.rgb = RGBColor(0xE6, 0xE9, 0xF7)
        yy += lh + gap
    return yy


def steps(s, x, y, w, items, gap=0.6):
    yy = y
    for i, (main, small) in enumerate(items):
        c = shape(s, MSO_SHAPE.OVAL, x, yy + 0.03, 0.32, 0.32)
        c.fill.solid(); c.fill.fore_color.rgb = VIOLET
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
        tb = s.shapes.add_textbox(Inches(x + 0.55), Inches(yy), Inches(w - 0.55), Inches(gap))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = main
        r.font.size = Pt(14.5); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = INK
        if small:
            p2 = tf.add_paragraph(); p2.space_before = Pt(1)
            r2 = p2.add_run(); r2.text = small
            r2.font.size = Pt(11.5); r2.font.name = FONT; r2.font.color.rgb = MUTED
        yy += gap
    return yy


def quote(s, x, y, w, h, before, hl=None, after=""):
    q = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        q.adjustments[0] = 0.08
    except Exception:
        pass
    q.fill.solid(); q.fill.fore_color.rgb = RGBColor(0x4F, 0x5B, 0xD8)
    tf = q.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.24); tf.margin_right = Inches(0.24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = before
    r.font.size = Pt(19); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
    if hl:
        r = p.add_run(); r.text = hl
        r.font.size = Pt(19); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = GOLD
    if after:
        r = p.add_run(); r.text = after
        r.font.size = Pt(19); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = WHITE
    return y + h


def box(s, x, y, w, text):
    b = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.62)
    try:
        b.adjustments[0] = 0.12
    except Exception:
        pass
    b.fill.solid(); b.fill.fore_color.rgb = PANEL
    b.line.color.rgb = LINE; b.line.width = Pt(0.75)
    rect(s, x, y + 0.08, 0.07, 0.46, GOLD)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.28); tf.margin_right = Inches(0.14)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.name = FONT; r.font.color.rgb = MUTED
    return y + 0.62


def code_block(s, x, y, w, h, code):
    c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        c.adjustments[0] = 0.05
    except Exception:
        pass
    c.fill.solid(); c.fill.fore_color.rgb = CODEBG
    c.line.color.rgb = LINE; c.line.width = Pt(0.75)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.1)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = ln if ln else " "
        r.font.size = Pt(11.5); r.font.name = MONO; r.font.color.rgb = CODE
    return y + h


def table(s, x, y, w, headers, rows, col_w, row_h=0.34, hdr_h=0.36):
    tbl = s.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y),
                             Inches(w), Inches(hdr_h + row_h * len(rows))).table
    for c, cw in enumerate(col_w):
        tbl.columns[c].width = Inches(cw)
    tbl.rows[0].height = Inches(hdr_h)
    for r in range(len(rows)):
        tbl.rows[r + 1].height = Inches(row_h)
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(11); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = GOLD
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = val
            r.font.size = Pt(11.5); r.font.name = FONT
            r.font.color.rgb = INK if j == 0 else RGBColor(0xE6, 0xE9, 0xF7)
            if j == 0:
                r.font.bold = True
    return y + hdr_h + row_h * len(rows)


def footer(s, n, total):
    lb = s.shapes.add_textbox(Inches(MX), Inches(FOOT_Y), Inches(6), Inches(0.3))
    p = lb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = DECK_TITLE
    r.font.size = Pt(9); r.font.name = FONT; r.font.color.rgb = RGBColor(0x66, 0x6E, 0x99)
    rb = s.shapes.add_textbox(Inches(W - 2.0), Inches(FOOT_Y), Inches(1.35), Inches(0.3))
    p2 = rb.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = "%d / %d" % (n, total)
    r2.font.size = Pt(9); r2.font.name = FONT; r2.font.color.rgb = RGBColor(0x66, 0x6E, 0x99)


def center_slide(kicker, title, sub=None, muted=None):
    s = new_slide()
    if kicker:
        k = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(W - 1.2), Inches(0.4))
        put_text(k.text_frame, kicker.upper(), 12, True, GREEN, align=PP_ALIGN.CENTER)
    t = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(W - 1.2), Inches(1.1))
    put_text(t.text_frame, title, 40, True, INK, align=PP_ALIGN.CENTER)
    yy = 3.7
    if sub:
        tb = s.shapes.add_textbox(Inches(1.4), Inches(yy), Inches(W - 2.8), Inches(0.9))
        put_text(tb.text_frame, sub, 16, False, RGBColor(0xC6, 0xCC, 0xE8), align=PP_ALIGN.CENTER)
        yy += 0.98
    if muted:
        mb = s.shapes.add_textbox(Inches(1.4), Inches(yy), Inches(W - 2.8), Inches(0.7))
        put_text(mb.text_frame, muted, 13, False, MUTED, align=PP_ALIGN.CENTER)
    return s, yy


def quiz(s, question, options, correct_idx, correct_note):
    tag_and_title(s, "Quick check", "Quick question", 30)
    tb = s.shapes.add_textbox(Inches(MX), Inches(CT), Inches(CW), Inches(0.6))
    put_text(tb.text_frame, question, 17, False, INK)
    yy = CT + 0.62
    for i, opt in enumerate(options):
        o = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MX, yy, CW, 0.55)
        try:
            o.adjustments[0] = 0.14
        except Exception:
            pass
        o.fill.solid()
        o.fill.fore_color.rgb = RGBColor(0x1E, 0x40, 0x34) if i == correct_idx else PANEL
        o.line.color.rgb = GREEN if i == correct_idx else LINE
        o.line.width = Pt(1.5 if i == correct_idx else 0.75)
        otf = o.text_frame; otf.word_wrap = True
        otf.margin_left = Inches(0.18); otf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = otf.paragraphs[0]
        r = p.add_run(); r.text = ("✓  " if i == correct_idx else "") + opt
        r.font.size = Pt(14); r.font.name = FONT
        r.font.color.rgb = GREEN if i == correct_idx else INK
        yy += 0.63
    fb = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, MX, yy + 0.08, CW, 0.62)
    try:
        fb.adjustments[0] = 0.1
    except Exception:
        pass
    fb.fill.solid(); fb.fill.fore_color.rgb = PANEL2
    fb.line.fill.background()
    tfb = fb.text_frame; tfb.word_wrap = True
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pfb = tfb.paragraphs[0]
    rfb = pfb.add_run(); rfb.text = "Correct: " + correct_note
    rfb.font.size = Pt(13); rfb.font.name = FONT; rfb.font.color.rgb = MUTED


# --- SLIDES data: one entry per HTML slide, same order, same copy -----------

DECK_NAME = "Deck Title"
SLIDES = [
    {"kind": "title", "tag": "Deck title", "title": "One line about the talk.", "chips": ["Topic one", "Topic two"]},
    {"kind": "content", "tag": "Label", "title": "Slide with a list", "items": ["One short point", "Another short point"]},
    {"kind": "content", "tag": "Label", "title": "Another content slide", "items": ["Point one", "Point two"]},
]


def slide_title(s, n, total, slide):
    k = s.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(W - 1.2), Inches(0.4))
    put_text(k.text_frame, slide["tag"], 13, True, GOLD, align=PP_ALIGN.CENTER)
    t = s.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(W - 1.2), Inches(1.0))
    put_text(t.text_frame, slide["title"], 26, True, INK, align=PP_ALIGN.CENTER)
    tags = slide.get("chips", [])
    iw = 0.9
    x0 = (W - iw * len(tags)) / 2
    for i, tg in enumerate(tags):
        c = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x0 + i * iw, 4.0, iw - 0.12, 0.34)
        try:
            c.adjustments[0] = 0.5
        except Exception:
            pass
        c.fill.solid(); c.fill.fore_color.rgb = PANEL
        c.line.color.rgb = LINE; c.line.width = Pt(0.75)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = tg
        r.font.size = Pt(11); r.font.name = FONT; r.font.color.rgb = MUTED
    footer(s, n, total)


def slide_content(s, n, total, slide):
    tag_and_title(s, slide["tag"], slide["title"])
    yy = CT + 0.25
    bullets(s, MX, yy, CW, slide["items"])
    footer(s, n, total)


def build_slide(s, n, total, slide):
    kind = slide.get("kind", "content")
    if kind == "title":
        slide_title(s, n, total, slide)
    else:
        slide_content(s, n, total, slide)


N = len(SLIDES)
for n, slide in enumerate(SLIDES, start=1):
    s = new_slide()
    build_slide(s, n, N, slide)

OUT = "".join(c for c in DECK_NAME if c.isalnum() or c in " -_").strip().replace(" ", "-") or "deck"
OUT = OUT + ".pptx"
prs.save(OUT)
print("saved", OUT, "with", N, "slides")